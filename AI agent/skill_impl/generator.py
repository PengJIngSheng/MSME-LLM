"""Local document-generation skills for the agent composer.

The module intentionally stays independent from Google Workspace connectors.
It creates local files, stores them in GridFS, and returns downloadable URLs.
"""

from __future__ import annotations

import csv
import io
import os
import re
import uuid
from dataclasses import dataclass
from html import escape as html_escape
from typing import Iterable

from docx import Document
from docx.shared import Inches as DocxInches
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


SUPPORTED_SKILLS = {"docx", "pdf", "pptx", "xlsx"}

CONTENT_TYPES = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf": "application/pdf",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


@dataclass(frozen=True)
class SkillResult:
    file_name: str
    file_url: str
    file_type: str
    content_type: str
    title: str


class SkillGenerationError(ValueError):
    """Raised when an unsupported or impossible skill request is received."""


def build_authoring_messages(
    skill_type: str,
    instruction: str,
    messages: list | None = None,
    attachments: list | None = None,
) -> list[dict]:
    """Build a compact LLM prompt that returns clean Markdown for the file writer."""
    skill_type = _normalize_skill_type(skill_type)
    recent_context = _latest_context(messages or [], max_chars=7000)
    attachment_context = _attachment_context(attachments or [])
    format_hint = {
        "docx": "Use a clear document structure with # title, ## sections, paragraphs, bullet lists, and tables only when useful.",
        "pdf": "Use polished report-style Markdown with # title, concise sections, bullets, and tables where helpful.",
        "pptx": "Use # as the deck title. Use ## for slide titles and short bullets below each slide. Keep each slide concise.",
        "xlsx": "Prefer one Markdown table with useful headers and rows. If the request is not naturally tabular, create a practical two-column table.",
    }[skill_type]

    system = (
        "You are a local file-authoring assistant. Convert the user's instruction and any provided context "
        f"into content for a .{skill_type} file.\n"
        "Return clean Markdown only. Do not use code fences. Do not mention that you are an AI. "
        "Do not include tool-call JSON or explanations outside the document content.\n"
        f"Formatting requirement: {format_hint}"
    )
    user = (
        f"User instruction:\n{(instruction or '').strip()}\n\n"
        f"Recent conversation context, if relevant:\n{recent_context or '(none)'}\n\n"
        f"Uploaded/generated file context:\n{attachment_context or '(none)'}"
    )
    return [{"role": "system", "content": system}, {"role": "user", "content": user}]


def generate_skill_file(
    skill_type: str,
    instruction: str,
    messages: list | None,
    fs,
    prepared_content: str | None = None,
    attachments: list | None = None,
) -> SkillResult:
    """Generate a file, write it into GridFS, and return its public metadata."""
    skill_type = _normalize_skill_type(skill_type)
    content = _clean_model_output(prepared_content or "")
    if len(content.strip()) < 20:
        content = _fallback_content(instruction, messages or [], attachments or [])

    title = _extract_title(content, instruction, skill_type)
    requested_file_stem = extract_requested_file_stem(instruction, skill_type)
    if skill_type == "docx":
        data = _build_docx(content, title)
    elif skill_type == "pdf":
        data = _build_pdf(content, title)
    elif skill_type == "pptx":
        data = _build_pptx(content, title)
    elif skill_type == "xlsx":
        data = _build_xlsx(content, title)
    else:
        raise SkillGenerationError(f"Unsupported skill type: {skill_type}")

    file_name = _safe_filename(requested_file_stem or title, skill_type, exact=bool(requested_file_stem), fs=fs)
    content_type = CONTENT_TYPES[skill_type]
    fs.put(
        data,
        filename=file_name,
        content_type=content_type,
        metadata={
            "source": "agent_skill",
            "skill_type": skill_type,
            "title": title,
            "requested_file_stem": requested_file_stem,
        },
    )
    return SkillResult(
        file_name=file_name,
        file_url=f"/uploads/{file_name}",
        file_type=skill_type,
        content_type=content_type,
        title=title,
    )


def _normalize_skill_type(skill_type: str) -> str:
    value = (skill_type or "").strip().lower().lstrip(".")
    if value not in SUPPORTED_SKILLS:
        raise SkillGenerationError(f"Unsupported skill type: {skill_type}")
    return value


def _clean_model_output(text: str) -> str:
    text = re.sub(r"<think>.*?</think>", "", text or "", flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"^```(?:markdown|md|text)?\s*|\s*```$", "", text.strip(), flags=re.IGNORECASE)
    return text.strip()


def _fallback_content(instruction: str, messages: list, attachments: list) -> str:
    context = _latest_context(messages, max_chars=5000)
    attachment_context = _attachment_context(attachments)
    parts = ["# Generated Document", "", "## Request", (instruction or "").strip()]
    if context:
        parts.extend(["", "## Relevant Context", context])
    if attachment_context:
        parts.extend(["", "## Files", attachment_context])
    return "\n".join(part for part in parts if part is not None)


def _latest_context(messages: list, max_chars: int = 6000) -> str:
    useful = []
    for msg in reversed(messages or []):
        if not isinstance(msg, dict):
            continue
        role = msg.get("role")
        content = _clean_model_output(str(msg.get("content", "") or "")).strip()
        if not content:
            continue
        if role == "assistant":
            if msg.get("pdf_url") or msg.get("generated_image_url"):
                continue
            if len(content) < 40:
                continue
            useful.append(f"Assistant:\n{content}")
        elif role == "user" and len(content) > 10:
            useful.append(f"User:\n{content}")
        if sum(len(item) for item in useful) >= max_chars:
            break
    return "\n\n".join(reversed(useful))[:max_chars].strip()


def _attachment_context(attachments: list | None) -> str:
    if not attachments:
        return ""
    lines = []
    for att in attachments:
        if not isinstance(att, dict):
            continue
        name = att.get("original_name") or att.get("saved_path") or att.get("file_id") or "file"
        size = att.get("size")
        content_type = att.get("content_type") or "unknown"
        size_note = f", {size} bytes" if size else ""
        lines.append(f"- {name} ({content_type}{size_note})")
    return "\n".join(lines)


def _extract_title(content: str, instruction: str, skill_type: str) -> str:
    for line in _iter_content_lines(content):
        heading = re.match(r"^#{1,3}\s+(.+)$", line)
        if heading:
            return _strip_inline_markdown(heading.group(1))[:80] or f"Generated {skill_type.upper()} File"
    for line in _iter_content_lines(instruction or content):
        plain = _strip_inline_markdown(line)
        if len(plain) > 4:
            return plain[:80]
    return f"Generated {skill_type.upper()} File"


def extract_requested_file_stem(instruction: str, skill_type: str = "") -> str:
    """Extract an explicit user-requested file name from the original prompt."""
    text = (instruction or "").strip()
    if not text:
        return ""

    file_words = (
        r"file|filename|document|doc|docx|pdf|spreadsheet|excel|xlsx|presentation|"
        r"slides?|slide\s+deck|deck|powerpoint|pptx|ppt|workbook|report"
    )
    patterns = [
        rf"(?is)\b(?:name|call|title)\s+it\s+(.+?)(?=$|[\n.!?。！？])",
        rf"(?is)\brename\s+(?:the\s+)?(?:(?:{file_words})\s+)?(?:file\s+)?(?:to|as)\s+(.+?)(?=$|[\n.!?。！？])",
        rf"(?is)(?:^|[\n.;!?])\s*(?:name|call|title)\s+(?:(?:it|this|the)\s+)?(?:{file_words})?\s*(?:as|to)?\s+(.+?)(?=$|[\n.!?。！？])",
        rf"(?is)(?:file\s*name|filename)\s*(?:is|as|to|should\s+be|:|=)?\s+(.+?)(?=$|[\n.!?。！？])",
        rf"(?is)(?:save\s+(?:it|the\s+(?:{file_words}))?\s+as)\s+(.+?)(?=$|[\n.!?。！？])",
        r"(?is)(?:命名为|命名成|取名为|文件名(?:叫|是|为)?|名字(?:叫|是|为)?|叫做|保存为|另存为)\s*[：:]?\s*(.+?)(?=$|[\n。！？])",
    ]
    for pattern in patterns:
        match = re.search(pattern, text)
        if not match:
            continue
        candidate = _clean_requested_file_stem(match.group(1), skill_type)
        if candidate:
            return candidate
    return ""


def _clean_requested_file_stem(candidate: str, skill_type: str = "") -> str:
    value = _strip_inline_markdown(candidate or "")
    value = value.strip().strip(" \"'`“”‘’")
    value = re.sub(r"\s+", " ", value)
    value = re.sub(r"^(?:it|this|the\s+file|the\s+document)\s+", "", value, flags=re.IGNORECASE).strip()
    value = re.sub(r"\b(?:please|thanks|thank you)\b.*$", "", value, flags=re.IGNORECASE).strip()
    value = value.strip(" .。")
    value = re.sub(r"\.(docx|pdf|pptx|ppt|xlsx|xls)$", "", value, flags=re.IGNORECASE).strip()
    if not value or len(value) > 90:
        return ""
    # Reject obvious accidental captures of the whole generation request.
    if re.search(r"\b(create|generate|make|include|columns?|rows?)\b", value, flags=re.IGNORECASE) and len(value.split()) > 6:
        return ""
    return value


def _filename_slug(stem: str, lowercase: bool = False) -> str:
    value = _strip_inline_markdown(stem or "")
    value = value.lower() if lowercase else value
    value = re.sub(r"[<>:\"/\\|?*\x00-\x1f]+", " ", value)
    value = re.sub(r"\s+", "-", value.strip())
    value = re.sub(r"[^\w.-]+", "-", value, flags=re.UNICODE)
    value = re.sub(r"-{2,}", "-", value).strip("-._")
    return value or "generated-file"


def _safe_filename(title: str, ext: str, exact: bool = False, fs=None) -> str:
    slug = _filename_slug(title, lowercase=not exact)
    if exact:
        base = slug[:64].strip("-._") or "generated-file"
        filename = f"{base}.{ext}"
        if fs is None or not fs.find_one({"filename": filename}):
            return filename
        return f"{base[:51].strip('-._')}-{uuid.uuid4().hex[:12]}.{ext}"
    return f"{slug[:48].strip('-._') or 'generated-file'}-{uuid.uuid4().hex[:12]}.{ext}"


def _iter_content_lines(text: str) -> Iterable[str]:
    for line in (text or "").replace("\r\n", "\n").split("\n"):
        line = line.strip()
        if line:
            yield line


def _strip_inline_markdown(text: str) -> str:
    text = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", text or "")
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"^\s*[-*]\s+", "", text)
    text = re.sub(r"[*_`>#]+", "", text)
    return re.sub(r"\s+", " ", text).strip()


def _parse_markdown_table(text: str) -> list[list[str]]:
    rows = []
    for line in (text or "").splitlines():
        stripped = line.strip()
        if not stripped.startswith("|") or "|" not in stripped[1:]:
            if rows:
                break
            continue
        cells = [cell.strip() for cell in stripped.strip("|").split("|")]
        if cells and all(re.fullmatch(r":?-{3,}:?", cell or "") for cell in cells):
            continue
        rows.append([_strip_inline_markdown(cell) for cell in cells])
    return rows if len(rows) >= 2 else []


def _parse_delimited_rows(text: str) -> list[list[str]]:
    lines = [line for line in (text or "").splitlines() if line.strip()]
    if len(lines) < 2:
        return []
    delimiter = "\t" if any("\t" in line for line in lines[:4]) else ","
    if delimiter == "," and not any("," in line for line in lines[:4]):
        return []
    try:
        reader = csv.reader(lines, delimiter=delimiter)
        rows = [[_strip_inline_markdown(cell.strip()) for cell in row] for row in reader]
        rows = [row for row in rows if any(row)]
    except csv.Error:
        return []
    return rows if len(rows) >= 2 else []


def _content_without_table(text: str) -> str:
    output = []
    in_table = False
    for line in (text or "").splitlines():
        stripped = line.strip()
        if stripped.startswith("|") and "|" in stripped[1:]:
            in_table = True
            continue
        if in_table and stripped:
            in_table = False
        if not in_table:
            output.append(line)
    return "\n".join(output).strip()


def _build_docx(content: str, title: str) -> bytes:
    document = Document()
    section = document.sections[0]
    section.top_margin = DocxInches(0.7)
    section.bottom_margin = DocxInches(0.7)
    section.left_margin = DocxInches(0.75)
    section.right_margin = DocxInches(0.75)
    document.add_heading(title, 0)

    table_rows = _parse_markdown_table(content)
    body = _content_without_table(content) if table_rows else content
    for raw in body.splitlines():
        line = raw.strip()
        if not line:
            continue
        heading = re.match(r"^(#{1,4})\s+(.+)$", line)
        bullet = re.match(r"^[-*]\s+(.+)$", line)
        numbered = re.match(r"^\d+[.)]\s+(.+)$", line)
        if heading:
            level = min(len(heading.group(1)), 3)
            document.add_heading(_strip_inline_markdown(heading.group(2)), level)
        elif bullet:
            document.add_paragraph(_strip_inline_markdown(bullet.group(1)), style="List Bullet")
        elif numbered:
            document.add_paragraph(_strip_inline_markdown(numbered.group(1)), style="List Number")
        else:
            document.add_paragraph(_strip_inline_markdown(line))

    if table_rows:
        document.add_paragraph()
        table = document.add_table(rows=1, cols=len(table_rows[0]))
        table.style = "Table Grid"
        for idx, cell in enumerate(table.rows[0].cells):
            cell.text = table_rows[0][idx]
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        for row in table_rows[1:]:
            cells = table.add_row().cells
            for idx, value in enumerate(row[: len(cells)]):
                cells[idx].text = value

    out = io.BytesIO()
    document.save(out)
    return out.getvalue()


def _build_xlsx(content: str, title: str) -> bytes:
    rows = _parse_markdown_table(content) or _parse_delimited_rows(content)
    if not rows:
        rows = [["Section", "Content"]]
        current_section = "Summary"
        for line in _iter_content_lines(content):
            heading = re.match(r"^#{1,4}\s+(.+)$", line)
            bullet = re.match(r"^[-*]\s+(.+)$", line)
            if heading:
                current_section = _strip_inline_markdown(heading.group(1))
            elif bullet:
                rows.append([current_section, _strip_inline_markdown(bullet.group(1))])
            else:
                rows.append([current_section, _strip_inline_markdown(line)])
        if len(rows) == 1:
            rows.append(["Request", title])

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    for row in rows:
        ws.append(row)

    header_fill = PatternFill("solid", fgColor="EEF2FF")
    for cell in ws[1]:
        cell.font = Font(bold=True, color="111827")
        cell.fill = header_fill
        cell.alignment = Alignment(vertical="center", wrap_text=True)
    ws.freeze_panes = "A2"

    for row in ws.iter_rows(min_row=2):
        for cell in row:
            cell.alignment = Alignment(vertical="top", wrap_text=True)

    for col_idx in range(1, ws.max_column + 1):
        max_len = 10
        for cell in ws[get_column_letter(col_idx)]:
            max_len = max(max_len, len(str(cell.value or "")))
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 3, 42)

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def _content_sections_for_slides(content: str, title: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    current_title = ""
    current_lines: list[str] = []

    for line in _iter_content_lines(content):
        heading = re.match(r"^##+\s+(.+)$", line)
        if heading:
            if current_title or current_lines:
                sections.append((current_title or title, current_lines))
            current_title = _strip_inline_markdown(heading.group(1))
            current_lines = []
            continue
        if line.startswith("# "):
            continue
        bullet = re.match(r"^[-*]\s+(.+)$", line)
        current_lines.append(_strip_inline_markdown(bullet.group(1) if bullet else line))

    if current_title or current_lines:
        sections.append((current_title or title, current_lines))
    if not sections:
        sections = [(title, [_strip_inline_markdown(line) for line in _iter_content_lines(content)])]
    return sections[:12]


def _build_pptx(content: str, title: str) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme = _pptx_theme_for_title(title)

    def add_box(slide, left, top, width, height, color, rounded=False):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
            left,
            top,
            width,
            height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _pptx_rgb(color)
        try:
            shape.line.fill.background()
        except Exception:
            shape.line.color.rgb = _pptx_rgb(color)
        return shape

    def add_text(slide, text, left, top, width, height, size, color, bold=False):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.02)
        frame.margin_right = Inches(0.02)
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = _pptx_rgb(color)
        return box

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_box(title_slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, theme["bg"])
    add_box(title_slide, Inches(0.65), Inches(0.72), Inches(1.35), Inches(0.08), theme["accent"])
    add_text(title_slide, title[:90], Inches(0.65), Inches(1.55), Inches(8.9), Inches(1.65), 40, theme["title"], True)
    add_text(title_slide, "Generated by MSME.AI", Inches(0.72), Inches(4.12), Inches(3.8), Inches(0.4), 13, theme["muted"], True)
    add_box(title_slide, Inches(10.15), Inches(0.62), Inches(1.6), Inches(1.6), theme["accent2"], True)
    add_box(title_slide, Inches(11.05), Inches(1.68), Inches(1.2), Inches(1.2), theme["accent"], True)
    add_box(title_slide, Inches(8.92), Inches(5.78), Inches(3.65), Inches(0.14), theme["accent"])

    slide_no = 1
    for section_title, lines in _content_sections_for_slides(content, title):
        for chunk in _chunk(lines, 5):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_box(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, theme["bg"])
            add_box(slide, Inches(0.55), Inches(0.55), Inches(12.2), Inches(6.3), theme["card"], True)
            add_box(slide, Inches(0.55), Inches(0.55), Inches(0.12), Inches(6.3), theme["accent"])
            add_text(slide, f"{slide_no:02d}", Inches(11.28), Inches(0.78), Inches(0.7), Inches(0.28), 11, theme["muted"], True)
            add_text(slide, section_title[:80], Inches(1.08), Inches(0.96), Inches(9.4), Inches(0.7), 29, theme["title"], True)
            body_box = slide.shapes.add_textbox(Inches(1.12), Inches(2.08), Inches(10.2), Inches(3.55))
            body = body_box.text_frame
            body.clear()
            body.word_wrap = True
            body.margin_left = Inches(0.02)
            body.margin_right = Inches(0.02)
            for idx, line in enumerate(chunk):
                paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
                paragraph.text = f"• {line[:185]}"
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(21)
                paragraph.font.color.rgb = _pptx_rgb(theme["body"])
                paragraph.space_after = Pt(10)
            add_box(slide, Inches(1.1), Inches(6.03), Inches(1.18), Inches(0.08), theme["accent2"])
            slide_no += 1

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _pptx_rgb(hex_color: str) -> RGBColor:
    value = (hex_color or "000000").strip().lstrip("#")
    if len(value) != 6:
        value = "000000"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _pptx_theme_for_title(title: str) -> dict:
    themes = [
        {
            "bg": "F6F8FB",
            "card": "FFFFFF",
            "accent": "2563EB",
            "accent2": "06B6D4",
            "title": "101828",
            "body": "344054",
            "muted": "667085",
        },
        {
            "bg": "07111F",
            "card": "0F1F35",
            "accent": "37D6C4",
            "accent2": "7C3AED",
            "title": "FFFFFF",
            "body": "D7E3F4",
            "muted": "8EA2BD",
        },
        {
            "bg": "F7F5EF",
            "card": "FFFFFF",
            "accent": "E95D35",
            "accent2": "111827",
            "title": "111827",
            "body": "384152",
            "muted": "7B8190",
        },
    ]
    seed = sum(ord(ch) for ch in (title or "presentation"))
    return themes[seed % len(themes)]


def _chunk(items: list[str], size: int) -> Iterable[list[str]]:
    filtered = [item for item in items if item]
    if not filtered:
        filtered = ["Summary"]
    for idx in range(0, len(filtered), size):
        yield filtered[idx : idx + size]


def _pdf_font_name() -> str:
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        return "STSong-Light"
    except Exception:
        return "Helvetica"


def _build_pdf(content: str, title: str) -> bytes:
    out = io.BytesIO()
    doc = SimpleDocTemplate(
        out,
        pagesize=A4,
        rightMargin=0.72 * inch,
        leftMargin=0.72 * inch,
        topMargin=0.7 * inch,
        bottomMargin=0.7 * inch,
        title=title,
    )
    font_name = _pdf_font_name()
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "SkillTitle",
        parent=styles["Title"],
        fontName=font_name,
        fontSize=20,
        leading=26,
        spaceAfter=18,
        textColor=colors.HexColor("#111827"),
    )
    heading_style = ParagraphStyle(
        "SkillHeading",
        parent=styles["Heading2"],
        fontName=font_name,
        fontSize=14,
        leading=18,
        spaceBefore=12,
        spaceAfter=8,
        textColor=colors.HexColor("#1f2937"),
    )
    body_style = ParagraphStyle(
        "SkillBody",
        parent=styles["BodyText"],
        fontName=font_name,
        fontSize=10.5,
        leading=16,
        spaceAfter=7,
    )
    bullet_style = ParagraphStyle(
        "SkillBullet",
        parent=body_style,
        leftIndent=14,
        firstLineIndent=-8,
    )

    flow = [Paragraph(html_escape(title), title_style)]
    table_rows = _parse_markdown_table(content)
    body = _content_without_table(content) if table_rows else content

    for raw in body.splitlines():
        line = raw.strip()
        if not line:
            continue
        if re.fullmatch(r"---+", line):
            flow.append(PageBreak())
            continue
        heading = re.match(r"^#{1,4}\s+(.+)$", line)
        bullet = re.match(r"^[-*]\s+(.+)$", line)
        if heading:
            flow.append(Paragraph(html_escape(_strip_inline_markdown(heading.group(1))), heading_style))
        elif bullet:
            flow.append(Paragraph(f"• {html_escape(_strip_inline_markdown(bullet.group(1)))}", bullet_style))
        else:
            flow.append(Paragraph(html_escape(_strip_inline_markdown(line)), body_style))

    if table_rows:
        flow.append(Spacer(1, 8))
        table = Table(table_rows, repeatRows=1)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EEF2FF")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#111827")),
            ("FONTNAME", (0, 0), (-1, -1), font_name),
            ("FONTSIZE", (0, 0), (-1, -1), 8.5),
            ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#D1D5DB")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
        ]))
        flow.append(table)

    doc.build(flow)
    return out.getvalue()
